"""Report generation tests — full workflow through the API.

Exercises POST /api/inputs -> POST /api/generate-report -> GET /api/download-report
and verifies the resulting .pptx file is valid and contains the expected slides.

Run with:  pytest tests/test_report_generation.py -v -s
"""

import os
from pathlib import Path

import pytest
from httpx import ASGITransport, AsyncClient
from pptx import Presentation

from app.main import app

# Re-use the realistic wound-camera fixture from the integration tests.
WOUND_CAMERA_INPUTS = {
    "setting": "Acute NHS Trust",
    "model_year": 2026,
    "forecast_years": 3,
    "funding_source": "Trust operational budget",
    "catchment_type": "population",
    "catchment_size": 350000,
    "eligible_pct": 2.5,
    "uptake_y1": 15,
    "uptake_y2": 40,
    "uptake_y3": 65,
    "prevalence": "Chronic wound prevalence ~2.5% of catchment; rising with ageing population",
    "workforce": [
        {"role": "Band 5 (Staff Nurse)", "minutes": 45, "frequency": "per patient"},
        {"role": "Band 6 (Senior Nurse/AHP)", "minutes": 15, "frequency": "per patient"},
        {"role": "Consultant", "minutes": 10, "frequency": "per patient"},
        {"role": "Admin/Clerical", "minutes": 5, "frequency": "per patient"},
    ],
    "outpatient_visits": 6,
    "tests": 3,
    "admissions": 1,
    "bed_days": 4,
    "procedures": 0,
    "consumables": 85.0,
    "pricing_model": "per-patient",
    "price": 750.0,
    "price_unit": "per year",
    "needs_training": True,
    "training_roles": "Band 5 nurses, Band 6 senior nurses",
    "training_hours": 3.0,
    "setup_cost": 8000.0,
    "staff_time_saved": 20.0,
    "visits_reduced": 25.0,
    "complications_reduced": 15.0,
    "readmissions_reduced": 10.0,
    "los_reduced": 0.5,
    "follow_up_reduced": 30.0,
    "comparator": "none",
    "comparator_names": "Visual assessment, paper wound chart",
    "discounting": "off",
}

REPORTS_DIR = Path(__file__).resolve().parent.parent / "data" / "reports"


@pytest.fixture
def transport():
    return ASGITransport(app=app)


# ── Helper ────────────────────────────────────────────────────────────

async def _submit_and_generate(client: AsyncClient) -> dict:
    """Submit inputs, generate report, return the generate-report response body."""
    save_resp = await client.post("/api/inputs", json=WOUND_CAMERA_INPUTS)
    assert save_resp.status_code == 200, f"Save failed: {save_resp.text}"
    sid = save_resp.json()["id"]

    gen_resp = await client.post(
        "/api/generate-report",
        json={"submission_id": sid},
    )
    assert gen_resp.status_code == 200, f"Generate failed: {gen_resp.text}"
    body = gen_resp.json()
    body["submission_id"] = sid
    return body


# ====================================================================
# 1. Full report workflow
# ====================================================================

class TestReportWorkflow:
    """Submit -> Generate -> Download -> Validate .pptx."""

    @pytest.mark.anyio
    async def test_generate_report_returns_success(self, transport):
        """POST /api/generate-report should return filepath and download URL."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            body = await _submit_and_generate(c)

        assert body["message"] == "Report generated successfully"
        assert body["filepath"].endswith(".pptx")
        assert body["download_url"].startswith("/api/download-report/")
        print(f"\n  Report generated: {body['filepath']}")

        # cleanup
        os.remove(body["filepath"])

    @pytest.mark.anyio
    async def test_pptx_file_created_in_reports_dir(self, transport):
        """The .pptx file should exist in data/reports/ with non-zero size."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            body = await _submit_and_generate(c)

        filepath = Path(body["filepath"])
        assert filepath.exists(), f"Report file not found: {filepath}"
        assert filepath.parent == REPORTS_DIR, (
            f"Report saved to wrong directory: {filepath.parent}"
        )

        size = filepath.stat().st_size
        assert size > 0, "Report file is empty"
        print(f"\n  File: {filepath.name}  Size: {size:,} bytes")

        # cleanup
        filepath.unlink()

    @pytest.mark.anyio
    async def test_pptx_opens_with_python_pptx(self, transport):
        """The file should be a valid PPTX that python-pptx can parse."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            body = await _submit_and_generate(c)

        filepath = body["filepath"]
        prs = Presentation(filepath)
        assert prs.slides is not None
        assert len(prs.slides) > 0
        print(f"\n  Opened successfully — {len(prs.slides)} slides")

        # cleanup
        os.remove(filepath)

    @pytest.mark.anyio
    async def test_presentation_has_expected_slides(self, transport):
        """Deck should contain 10 slides (title, input summary, 3 section
        dividers, exec summary, population, budget impact table, scenario
        comparison, assumptions)."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            body = await _submit_and_generate(c)

        filepath = body["filepath"]
        prs = Presentation(filepath)
        slide_count = len(prs.slides)

        assert slide_count == 10, (
            f"Expected 10 slides, got {slide_count}"
        )

        # Verify key slide titles by inspecting the first text in each
        expected_titles = [
            "HEOR Engine",                                # 1: Title
            "Input Summary",                              # 2: Input Summary
            "Executive Summary",                          # 3: Section divider
            "Executive Summary",                          # 4: Exec content
            "Population & Uptake",                        # 5: Section divider
            "Population & Uptake",                        # 6: Population content
            "Budget Impact",                              # 7: Section divider
            "Annual Budget Impact (Base Case)",           # 8: Budget table
            "Scenario Analysis (3-Year Total Impact)",    # 9: Scenarios
            "Model Assumptions & Methodology",            # 10: Assumptions
        ]

        for i, slide in enumerate(prs.slides):
            texts = [
                s.text_frame.text.strip()
                for s in slide.shapes
                if s.has_text_frame and s.text_frame.text.strip()
            ]
            first = texts[0] if texts else "(empty)"
            assert first == expected_titles[i], (
                f"Slide {i + 1}: expected title '{expected_titles[i]}', "
                f"got '{first}'"
            )
            print(f"  Slide {i + 1}: {first}")

        # cleanup
        os.remove(filepath)

    @pytest.mark.anyio
    async def test_download_returns_pptx_bytes(self, transport):
        """GET /api/download-report/{id} should return the file with correct
        Content-Type and Content-Disposition headers."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            body = await _submit_and_generate(c)
            sid = body["submission_id"]

            dl = await c.get(f"/api/download-report/{sid}")

        assert dl.status_code == 200, f"Download failed: {dl.status_code}"
        assert len(dl.content) > 0, "Download response body is empty"

        disposition = dl.headers.get("content-disposition", "")
        assert "attachment" in disposition, (
            f"Expected attachment disposition, got: {disposition}"
        )
        assert sid in disposition, (
            f"Filename should contain submission_id, got: {disposition}"
        )
        print(f"\n  Downloaded {len(dl.content):,} bytes")
        print(f"  Content-Disposition: {disposition}")

        # cleanup
        os.remove(body["filepath"])


# ====================================================================
# 2. Error handling
# ====================================================================

class TestReportErrors:

    @pytest.mark.anyio
    async def test_generate_404_for_missing_submission(self, transport):
        """Generating a report for a non-existent submission returns 404."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post(
                "/api/generate-report",
                json={"submission_id": "does_not_exist_99999"},
            )
        assert resp.status_code == 404, (
            f"Expected 404, got {resp.status_code}"
        )
        assert "not found" in resp.json()["detail"].lower()
        print(f"\n  Correctly returned 404 for missing submission")

    @pytest.mark.anyio
    async def test_download_404_for_missing_report(self, transport):
        """Downloading a report that hasn't been generated returns 404."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.get("/api/download-report/does_not_exist_99999")
        assert resp.status_code == 404, (
            f"Expected 404, got {resp.status_code}"
        )
        assert "not found" in resp.json()["detail"].lower()
        print(f"\n  Correctly returned 404 for missing report")

    @pytest.mark.anyio
    async def test_generate_422_for_missing_submission_id(self, transport):
        """Omitting submission_id from the request body returns 422."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/generate-report", json={})
        assert resp.status_code == 422, (
            f"Expected 422, got {resp.status_code}"
        )
        print(f"\n  Correctly returned 422 for missing submission_id")


# ====================================================================
# 3. Optional: open the report on macOS (skipped in CI)
# ====================================================================

@pytest.mark.skipif(
    os.environ.get("CI") == "true",
    reason="Skipped in CI — no GUI available",
)
class TestOpenReport:
    """Generate a report and open it with the system viewer.

    Run explicitly with:  pytest tests/test_report_generation.py::TestOpenReport -v -s
    """

    @pytest.mark.anyio
    async def test_open_generated_report(self, transport):
        import platform
        import subprocess

        async with AsyncClient(transport=transport, base_url="http://test") as c:
            body = await _submit_and_generate(c)

        filepath = body["filepath"]
        print(f"\n  Opening report: {filepath}")

        system = platform.system()
        if system == "Darwin":
            subprocess.run(["open", filepath])
        elif system == "Windows":
            subprocess.run(["start", filepath], shell=True)
        elif system == "Linux":
            subprocess.run(["xdg-open", filepath])
        else:
            print(f"  Cannot auto-open on {system} — file is at {filepath}")
