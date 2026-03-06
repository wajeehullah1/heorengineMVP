"""Markov integration tests — full API workflow for ICER, BIA→CEA, combined report.

Exercises the Markov/ICER endpoints end-to-end through the FastAPI app:
    POST /api/calculate-icer
    POST /api/calculate-icer-from-bia
    POST /api/generate-combined-report
    GET  /api/download-combined-report/{id}

Also covers validation and error handling for invalid inputs.

Run with:  pytest tests/test_markov_integration.py -v -s
"""

import os
from pathlib import Path

import pytest
from httpx import ASGITransport, AsyncClient
from pptx import Presentation

from app.main import app
from engines.markov.runner import check_r_installed


# ── Skip entire module if R is not installed ─────────────────────────

def pytest_configure():
    if not check_r_installed():
        pytest.skip(
            "R not installed — skipping Markov integration tests.\n"
            "Install R from https://cloud.r-project.org/",
            allow_module_level=True,
        )


# ── Shared fixtures ──────────────────────────────────────────────────

WOUND_CAMERA_INPUTS = {
    "setting": "Acute NHS Trust",
    "model_year": 2026,
    "forecast_years": 3,
    "funding_source": "Trust operational budget",
    "catchment_type": "population",
    "catchment_size": 350000,
    "eligible_pct": 2.5,
    "uptake_y1": 15,
    "uptake_y2": 40,
    "uptake_y3": 65,
    "prevalence": "Chronic wound prevalence ~2.5%",
    "workforce": [
        {"role": "Band 5 (Staff Nurse)", "minutes": 45, "frequency": "per patient"},
        {"role": "Band 6 (Senior Nurse/AHP)", "minutes": 15, "frequency": "per patient"},
        {"role": "Consultant", "minutes": 10, "frequency": "per patient"},
        {"role": "Admin/Clerical", "minutes": 5, "frequency": "per patient"},
    ],
    "outpatient_visits": 6,
    "tests": 3,
    "admissions": 1,
    "bed_days": 4,
    "procedures": 0,
    "consumables": 85.0,
    "pricing_model": "per-patient",
    "price": 750.0,
    "price_unit": "per year",
    "needs_training": True,
    "training_roles": "Band 5 nurses, Band 6 senior nurses",
    "training_hours": 3.0,
    "setup_cost": 8000.0,
    "staff_time_saved": 20.0,
    "visits_reduced": 25.0,
    "complications_reduced": 15.0,
    "readmissions_reduced": 10.0,
    "los_reduced": 0.5,
    "follow_up_reduced": 30.0,
    "comparator": "none",
    "comparator_names": "Visual assessment, paper wound chart",
    "discounting": "off",
}

SAMPLE_MARKOV_INPUTS = {
    "intervention_name": "Wound Camera",
    "time_horizon": 5,
    "cycle_length": 1.0,
    "discount_rate": 0.035,
    "prob_death_standard": 0.08,
    "cost_standard_annual": 8000.0,
    "utility_standard": 0.65,
    "prob_death_treatment": 0.04,
    "cost_treatment_annual": 9500.0,
    "cost_treatment_initial": 2000.0,
    "utility_treatment": 0.78,
}

REPORTS_DIR = Path(__file__).resolve().parent.parent / "data" / "reports"


@pytest.fixture
def transport():
    return ASGITransport(app=app)


async def _create_bia_submission(client: AsyncClient) -> str:
    """Submit wound camera inputs and return the submission ID."""
    resp = await client.post("/api/inputs", json=WOUND_CAMERA_INPUTS)
    assert resp.status_code == 200, f"BIA save failed: {resp.text}"
    return resp.json()["id"]


# ====================================================================
# 1. Basic ICER API endpoint
# ====================================================================

class TestICEREndpoint:
    """POST /api/calculate-icer with full MarkovInputs."""

    @pytest.mark.anyio
    async def test_returns_valid_markov_results(self, transport):
        """Response should contain all MarkovResults fields."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/calculate-icer", json=SAMPLE_MARKOV_INPUTS)

        assert resp.status_code == 200, f"ICER failed: {resp.text}"
        body = resp.json()

        # Verify top-level structure
        for key in ("standard_care", "treatment", "incremental_cost",
                     "incremental_qalys", "icer", "interpretation",
                     "cost_effective_25k", "cost_effective_35k"):
            assert key in body, f"Missing key '{key}'"

        # Verify arm structure
        for arm in ("standard_care", "treatment"):
            assert "total_cost" in body[arm], f"{arm} missing total_cost"
            assert "total_qalys" in body[arm], f"{arm} missing total_qalys"

        print(f"\n  Standard care: £{body['standard_care']['total_cost']:,.0f} / "
              f"{body['standard_care']['total_qalys']:.2f} QALYs")
        print(f"  Treatment:     £{body['treatment']['total_cost']:,.0f} / "
              f"{body['treatment']['total_qalys']:.2f} QALYs")
        print(f"  ICER: £{body['icer']:,.0f}/QALY")
        print(f"  Interpretation: {body['interpretation']}")

    @pytest.mark.anyio
    async def test_icer_calculation_is_correct(self, transport):
        """ICER should equal incremental_cost / incremental_qalys."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/calculate-icer", json=SAMPLE_MARKOV_INPUTS)

        body = resp.json()
        expected_icer = body["incremental_cost"] / body["incremental_qalys"]
        assert body["icer"] == pytest.approx(expected_icer, rel=0.01), (
            f"ICER {body['icer']} != incremental_cost/incremental_qalys = {expected_icer}"
        )

    @pytest.mark.anyio
    async def test_interpretation_matches_icer(self, transport):
        """Interpretation string should be consistent with ICER value."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/calculate-icer", json=SAMPLE_MARKOV_INPUTS)

        body = resp.json()
        icer = body["icer"]

        # With our sample inputs, treatment costs more but gains QALYs.
        # ICER should be positive and below £25k (cost-effective).
        assert icer > 0, "ICER should be positive for this scenario"
        assert icer < 25000, f"Expected ICER < £25k, got £{icer:,.0f}"
        assert "cost-effective" in body["interpretation"].lower(), (
            f"Expected 'cost-effective' in interpretation, got: {body['interpretation']}"
        )
        assert body["cost_effective_25k"] is True
        assert body["cost_effective_35k"] is True

    @pytest.mark.anyio
    async def test_treatment_gains_qalys(self, transport):
        """Treatment should produce more QALYs (lower mortality + higher utility)."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/calculate-icer", json=SAMPLE_MARKOV_INPUTS)

        body = resp.json()
        assert body["treatment"]["total_qalys"] > body["standard_care"]["total_qalys"]
        assert body["incremental_qalys"] > 0


# ====================================================================
# 2. BIA → Markov conversion
# ====================================================================

class TestBIAToMarkov:
    """POST /api/calculate-icer-from-bia — derive CEA from BIA submission."""

    @pytest.mark.anyio
    async def test_derives_markov_params_from_bia(self, transport):
        """Should return combined BIA summary + Markov results."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            resp = await c.post("/api/calculate-icer-from-bia", json={
                "submission_id": sid,
                "mortality_reduction": 50,
                "utility_gain": 0.13,
                "base_mortality": 0.08,
                "base_utility": 0.65,
                "time_horizon": 5,
            })

        assert resp.status_code == 200, f"BIA→ICER failed: {resp.text}"
        body = resp.json()

        # Verify response structure
        assert body["submission_id"] == sid
        assert "bia_summary" in body
        assert "markov_inputs" in body
        assert "markov_results" in body

        # Verify BIA summary fields
        bia = body["bia_summary"]
        assert bia["eligible_patients"] > 0
        assert len(bia["treated_patients"]) == 3
        assert len(bia["annual_budget_impact"]) == 3
        assert len(bia["cost_per_patient"]) == 3

        print(f"\n  Eligible patients: {bia['eligible_patients']}")
        print(f"  BIA cost/patient Y1: £{bia['cost_per_patient'][0]:,.2f}")

    @pytest.mark.anyio
    async def test_mortality_reduction_applied(self, transport):
        """Treatment mortality should be base * (1 - reduction/100)."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            resp = await c.post("/api/calculate-icer-from-bia", json={
                "submission_id": sid,
                "mortality_reduction": 50,
                "utility_gain": 0.10,
                "base_mortality": 0.10,
                "base_utility": 0.65,
            })

        body = resp.json()
        mi = body["markov_inputs"]

        # 50% reduction of 0.10 → 0.05
        assert mi["prob_death_treatment"] == pytest.approx(0.05), (
            f"Expected 0.05, got {mi['prob_death_treatment']}"
        )
        assert mi["prob_death_standard"] == pytest.approx(0.10)
        print(f"\n  Mortality: standard={mi['prob_death_standard']}, "
              f"treatment={mi['prob_death_treatment']}")

    @pytest.mark.anyio
    async def test_utility_gain_applied(self, transport):
        """Treatment utility should be base_utility + utility_gain."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            resp = await c.post("/api/calculate-icer-from-bia", json={
                "submission_id": sid,
                "mortality_reduction": 25,
                "utility_gain": 0.15,
                "base_mortality": 0.08,
                "base_utility": 0.65,
            })

        body = resp.json()
        mi = body["markov_inputs"]

        assert mi["utility_standard"] == pytest.approx(0.65)
        assert mi["utility_treatment"] == pytest.approx(0.80), (
            f"Expected 0.80, got {mi['utility_treatment']}"
        )
        print(f"\n  Utility: standard={mi['utility_standard']}, "
              f"treatment={mi['utility_treatment']}")

    @pytest.mark.anyio
    async def test_treatment_cost_matches_bia(self, transport):
        """Treatment annual cost should come from BIA cost_per_patient."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            # Get BIA results to know cost_per_patient
            bia_resp = await c.post(
                "/api/calculate-bia", params={"submission_id": sid}
            )
            bia_cpp = bia_resp.json()["base"]["cost_per_patient"][0]

            # Now run the conversion
            resp = await c.post("/api/calculate-icer-from-bia", json={
                "submission_id": sid,
                "mortality_reduction": 50,
                "utility_gain": 0.10,
            })

        body = resp.json()
        mi = body["markov_inputs"]

        assert mi["cost_treatment_annual"] == pytest.approx(bia_cpp, rel=0.01), (
            f"Treatment cost £{mi['cost_treatment_annual']:,.2f} should match "
            f"BIA cost/patient £{bia_cpp:,.2f}"
        )
        print(f"\n  BIA cost/patient: £{bia_cpp:,.2f}")
        print(f"  Markov treatment annual: £{mi['cost_treatment_annual']:,.2f}")

    @pytest.mark.anyio
    async def test_markov_results_present(self, transport):
        """The Markov results should contain a valid ICER."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            resp = await c.post("/api/calculate-icer-from-bia", json={
                "submission_id": sid,
                "mortality_reduction": 50,
                "utility_gain": 0.13,
            })

        body = resp.json()
        mr = body["markov_results"]
        assert "icer" in mr
        assert "interpretation" in mr
        assert mr["incremental_qalys"] > 0, "Treatment should gain QALYs"
        print(f"\n  ICER: {mr['icer']}")
        print(f"  Interpretation: {mr['interpretation']}")


# ====================================================================
# 3. Combined report generation
# ====================================================================

class TestCombinedReport:
    """POST /api/generate-combined-report → 16-slide BIA+CEA deck."""

    @pytest.mark.anyio
    async def test_generates_combined_report(self, transport):
        """Should return filepath, download URL, and summary."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            resp = await c.post("/api/generate-combined-report", json={
                "bia_submission_id": sid,
                "markov_params": SAMPLE_MARKOV_INPUTS,
            })

        assert resp.status_code == 200, f"Combined report failed: {resp.text}"
        body = resp.json()

        assert "filepath" in body
        assert body["download_url"] == f"/api/download-combined-report/{sid}"
        assert "summary" in body
        assert body["message"] == "Combined BIA + CEA report generated successfully"

        summary = body["summary"]
        assert "budget_impact_3yr" in summary
        assert "icer" in summary
        assert "recommendation" in summary
        assert isinstance(summary["cost_effective_25k"], bool)
        assert isinstance(summary["cost_effective_35k"], bool)

        print(f"\n  BIA 3-year: {summary['budget_impact_3yr']}")
        print(f"  ICER: {summary['icer']}")
        print(f"  Recommendation: {summary['recommendation']}")

        # cleanup
        if body["filepath"] and os.path.exists(body["filepath"]):
            os.remove(body["filepath"])

    @pytest.mark.anyio
    async def test_pptx_created_with_bia_and_cea_slides(self, transport):
        """PPTX should exist, be valid, and have 15+ slides."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            resp = await c.post("/api/generate-combined-report", json={
                "bia_submission_id": sid,
                "markov_params": SAMPLE_MARKOV_INPUTS,
            })

        body = resp.json()
        filepath = body["filepath"]
        assert filepath is not None
        assert os.path.exists(filepath), f"PPTX not found: {filepath}"

        # Open and verify slide count
        prs = Presentation(filepath)
        slide_count = len(prs.slides)
        assert slide_count >= 15, (
            f"Expected 15+ slides (BIA + CEA), got {slide_count}"
        )

        # Verify we have both BIA and CEA content by checking slide text
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        all_text.append(text)

        joined = " ".join(all_text).lower()
        assert "budget impact" in joined, "Should contain BIA slides"
        assert "cost-effectiveness" in joined, "Should contain CEA slides"

        print(f"\n  Slides: {slide_count}")
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    print(f"    Slide {i}: {shape.text_frame.text.strip()[:60]}")
                    break

        # cleanup
        os.remove(filepath)

    @pytest.mark.anyio
    async def test_download_combined_report(self, transport):
        """GET /api/download-combined-report should return the PPTX file."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            gen_resp = await c.post("/api/generate-combined-report", json={
                "bia_submission_id": sid,
                "markov_params": SAMPLE_MARKOV_INPUTS,
            })
            assert gen_resp.status_code == 200

            dl_resp = await c.get(f"/api/download-combined-report/{sid}")

        assert dl_resp.status_code == 200, f"Download failed: {dl_resp.status_code}"
        assert len(dl_resp.content) > 0, "Downloaded file is empty"

        disposition = dl_resp.headers.get("content-disposition", "")
        assert "attachment" in disposition
        assert sid in disposition

        print(f"\n  Downloaded {len(dl_resp.content):,} bytes")
        print(f"  Content-Disposition: {disposition}")

        # cleanup
        report_path = REPORTS_DIR / f"BIA_CEA_{sid}.pptx"
        if report_path.exists():
            report_path.unlink()

    @pytest.mark.anyio
    async def test_intervention_name_override(self, transport):
        """Optional intervention_name should override markov_params name."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            resp = await c.post("/api/generate-combined-report", json={
                "bia_submission_id": sid,
                "markov_params": SAMPLE_MARKOV_INPUTS,
                "intervention_name": "Custom Device Name",
            })

        assert resp.status_code == 200
        # cleanup
        body = resp.json()
        if body["filepath"] and os.path.exists(body["filepath"]):
            os.remove(body["filepath"])


# ====================================================================
# 4. Error cases
# ====================================================================

class TestICERValidationErrors:
    """Invalid inputs to /api/calculate-icer should return 422."""

    @pytest.mark.anyio
    async def test_probability_above_1(self, transport):
        """Probability > 1 should be rejected."""
        bad = {**SAMPLE_MARKOV_INPUTS, "prob_death_standard": 1.5}
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/calculate-icer", json=bad)

        assert resp.status_code == 422, (
            f"Expected 422 for prob > 1, got {resp.status_code}"
        )
        print(f"\n  Correctly rejected prob_death_standard=1.5")

    @pytest.mark.anyio
    async def test_probability_below_0(self, transport):
        """Negative probability should be rejected."""
        bad = {**SAMPLE_MARKOV_INPUTS, "prob_death_treatment": -0.1}
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/calculate-icer", json=bad)

        assert resp.status_code == 422, (
            f"Expected 422 for prob < 0, got {resp.status_code}"
        )
        print(f"\n  Correctly rejected prob_death_treatment=-0.1")

    @pytest.mark.anyio
    async def test_negative_cost(self, transport):
        """Negative annual cost should be rejected."""
        bad = {**SAMPLE_MARKOV_INPUTS, "cost_standard_annual": -5000}
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/calculate-icer", json=bad)

        assert resp.status_code == 422, (
            f"Expected 422 for negative cost, got {resp.status_code}"
        )
        print(f"\n  Correctly rejected cost_standard_annual=-5000")

    @pytest.mark.anyio
    async def test_missing_required_fields(self, transport):
        """Omitting required fields should return 422."""
        partial = {
            "intervention_name": "Incomplete",
            "time_horizon": 5,
        }
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/calculate-icer", json=partial)

        assert resp.status_code == 422, (
            f"Expected 422 for missing fields, got {resp.status_code}"
        )
        print(f"\n  Correctly rejected incomplete MarkovInputs")

    @pytest.mark.anyio
    async def test_utility_above_1(self, transport):
        """Utility weight > 1 should be rejected."""
        bad = {**SAMPLE_MARKOV_INPUTS, "utility_treatment": 1.5}
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/calculate-icer", json=bad)

        assert resp.status_code == 422, (
            f"Expected 422 for utility > 1, got {resp.status_code}"
        )
        print(f"\n  Correctly rejected utility_treatment=1.5")

    @pytest.mark.anyio
    async def test_empty_intervention_name(self, transport):
        """Empty intervention name should be rejected (min_length=1)."""
        bad = {**SAMPLE_MARKOV_INPUTS, "intervention_name": ""}
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/calculate-icer", json=bad)

        assert resp.status_code == 422, (
            f"Expected 422 for empty name, got {resp.status_code}"
        )
        print(f"\n  Correctly rejected empty intervention_name")


class TestBIAToMarkovErrors:
    """Error handling for /api/calculate-icer-from-bia."""

    @pytest.mark.anyio
    async def test_invalid_submission_id(self, transport):
        """Non-existent submission should return 404."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/calculate-icer-from-bia", json={
                "submission_id": "does_not_exist_99999",
                "mortality_reduction": 50,
                "utility_gain": 0.10,
            })

        assert resp.status_code == 404, (
            f"Expected 404, got {resp.status_code}"
        )
        assert "not found" in resp.json()["detail"].lower()
        print(f"\n  Correctly returned 404 for invalid submission")

    @pytest.mark.anyio
    async def test_missing_mortality_reduction(self, transport):
        """Omitting mortality_reduction should return 422."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            resp = await c.post("/api/calculate-icer-from-bia", json={
                "submission_id": sid,
                "utility_gain": 0.10,
            })

        assert resp.status_code == 422, (
            f"Expected 422 for missing mortality_reduction, got {resp.status_code}"
        )

    @pytest.mark.anyio
    async def test_mortality_reduction_out_of_range(self, transport):
        """mortality_reduction > 100 should be rejected."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            resp = await c.post("/api/calculate-icer-from-bia", json={
                "submission_id": sid,
                "mortality_reduction": 150,
                "utility_gain": 0.10,
            })

        assert resp.status_code == 422, (
            f"Expected 422 for mortality_reduction=150, got {resp.status_code}"
        )


class TestCombinedReportErrors:
    """Error handling for /api/generate-combined-report."""

    @pytest.mark.anyio
    async def test_invalid_bia_submission_id(self, transport):
        """Non-existent BIA submission should return 404."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/generate-combined-report", json={
                "bia_submission_id": "does_not_exist_99999",
                "markov_params": SAMPLE_MARKOV_INPUTS,
            })

        assert resp.status_code == 404, (
            f"Expected 404, got {resp.status_code}"
        )
        assert "not found" in resp.json()["detail"].lower()
        print(f"\n  Correctly returned 404 for invalid BIA submission")

    @pytest.mark.anyio
    async def test_invalid_markov_params(self, transport):
        """Invalid Markov params should return 422."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            sid = await _create_bia_submission(c)

            bad_markov = {**SAMPLE_MARKOV_INPUTS, "prob_death_standard": 2.0}
            resp = await c.post("/api/generate-combined-report", json={
                "bia_submission_id": sid,
                "markov_params": bad_markov,
            })

        assert resp.status_code == 422, (
            f"Expected 422 for invalid Markov params, got {resp.status_code}"
        )
        print(f"\n  Correctly returned 422 for invalid Markov params")

    @pytest.mark.anyio
    async def test_download_404_before_generation(self, transport):
        """Downloading before generating should return 404."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.get(
                "/api/download-combined-report/does_not_exist_99999"
            )

        assert resp.status_code == 404, (
            f"Expected 404, got {resp.status_code}"
        )
        print(f"\n  Correctly returned 404 for missing combined report")

    @pytest.mark.anyio
    async def test_missing_markov_params(self, transport):
        """Omitting markov_params should return 422."""
        async with AsyncClient(transport=transport, base_url="http://test") as c:
            resp = await c.post("/api/generate-combined-report", json={
                "bia_submission_id": "some_id",
            })

        assert resp.status_code == 422, (
            f"Expected 422 for missing markov_params, got {resp.status_code}"
        )
