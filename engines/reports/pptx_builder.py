"""PowerPoint report builder for BIA and CEA (Markov) results.

Generates branded PPTX slide decks from BIAInputs / MarkovInputs + results,
suitable for NICE submissions and NHS trust stakeholder meetings.
"""

from __future__ import annotations

from datetime import date
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

from engines.bia.schema import BIAInputs, BIAResults
from engines.markov.schema import MarkovInputs, MarkovResults

# ── Theme constants ───────────────────────────────────────────────────

NAVY = RGBColor(0x0D, 0x20, 0x44)
GREEN = RGBColor(0x2D, 0x6A, 0x4F)
LIGHT_GREEN = RGBColor(0x52, 0xB7, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x66, 0x66, 0x66)
LABEL_GREY = RGBColor(0x5A, 0x7F, 0xA8)
VALUE_WHITE = RGBColor(0xE8, 0xED, 0xF5)
RED = RGBColor(0xDC, 0x26, 0x26)
BAR_CONSERVATIVE = RGBColor(0x74, 0xC6, 0x9D)
BAR_BASE = RGBColor(0x52, 0xB7, 0x88)
BAR_OPTIMISTIC = RGBColor(0x2D, 0x6A, 0x4F)
AMBER = RGBColor(0xFF, 0xC1, 0x07)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)

FONT_NAME = "Calibri"
SLIDE_WIDTH = 13_333_400  # 16:9 widescreen EMU
SLIDE_HEIGHT = 7_500_938

REPORTS_DIR = Path(__file__).resolve().parent.parent.parent / "data" / "reports"


# ── Helpers ───────────────────────────────────────────────────────────

def set_heor_theme(prs: Presentation) -> None:
    """Set slide dimensions to 16:9 widescreen."""
    prs.slide_width = Emu(SLIDE_WIDTH)
    prs.slide_height = Emu(SLIDE_HEIGHT)


def _add_text(
    text_frame,
    text: str,
    size_pt: int = 14,
    bold: bool = False,
    color: RGBColor = BLACK,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """Add a formatted paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_NAME


def _add_shape_rect(slide, left, top, width, height, fill_color):
    """Add a filled rectangle shape to the slide."""
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height) -> object:
    """Add a textbox and return it. Clears the default paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    # Clear default empty paragraph text
    if txBox.text_frame.paragraphs:
        txBox.text_frame.paragraphs[0].text = ""
    return txBox


def _fmt_gbp(value: float) -> str:
    """Format a number as £ with commas, no decimals for large numbers."""
    if abs(value) >= 1000:
        return f"£{value:,.0f}"
    return f"£{value:,.2f}"


def _add_navy_title_bar(slide, title_text: str) -> None:
    """Add a navy title bar at the top of a content slide."""
    _add_shape_rect(
        slide,
        Inches(0), Inches(0),
        Emu(SLIDE_WIDTH), Inches(1.2),
        NAVY,
    )
    tb = _add_textbox(slide, Inches(0.5), Inches(0.17), Inches(12.3), Inches(0.86))
    tf = tb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = title_text
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME


def _add_table(slide, rows, cols, left, top, width, height):
    """Add a table shape and return the table object."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def _style_table_cell(cell, text, size_pt=11, bold=False, color=BLACK, bg=None, align=PP_ALIGN.LEFT):
    """Set cell text formatting."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_NAME
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


# ── Slide builders ────────────────────────────────────────────────────

def add_title_slide(prs: Presentation, intervention_name: str, report_date: str) -> None:
    """Slide 1: Title slide with HEOR Engine branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Navy background
    _add_shape_rect(slide, Inches(0), Inches(0), Emu(SLIDE_WIDTH), Emu(SLIDE_HEIGHT), NAVY)

    # Green accent bar
    _add_shape_rect(slide, Inches(0.6), Inches(1.5), Inches(1.5), Inches(0.06), LIGHT_GREEN)

    # "HEOR Engine" branding
    tb = _add_textbox(slide, Inches(0.6), Inches(0.6), Inches(8), Inches(0.8))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "HEOR Engine"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = LIGHT_GREEN
    run.font.name = FONT_NAME

    # Main title: intervention setting
    tb = _add_textbox(slide, Inches(0.6), Inches(1.8), Inches(11), Inches(1.5))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = intervention_name
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    # Subtitle
    tb = _add_textbox(slide, Inches(0.6), Inches(3.5), Inches(11), Inches(1))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "Budget Impact Analysis"
    run.font.size = Pt(22)
    run.font.color.rgb = LIGHT_GREEN
    run.font.name = FONT_NAME
    _add_text(tf, report_date, size_pt=16, color=WHITE, alignment=PP_ALIGN.LEFT)


def add_section_divider(prs: Presentation, section_title: str) -> None:
    """Section divider slide: green background, white centred title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Green background
    _add_shape_rect(slide, Inches(0), Inches(0), Emu(SLIDE_WIDTH), Emu(SLIDE_HEIGHT), GREEN)

    # Centred title
    tb = _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_title
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME


def add_input_summary_slide(prs: Presentation, inputs: BIAInputs) -> None:
    """Input Summary slide: two-column layout of key BIA inputs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_title_bar(slide, "Input Summary")

    # ── Left column items ─────────────────────────────────────────────
    left_items = [
        ("Setting", inputs.setting.value),
        ("Model timeframe", f"FY{inputs.model_year}, {inputs.forecast_years}-year horizon"),
        ("Funding route", inputs.funding_source.value),
        ("Catchment", f"{inputs.catchment_size:,} {inputs.catchment_type.value}"),
        ("Eligible patients", f"{inputs.eligible_patients:,}"),
    ]

    # ── Right column items ────────────────────────────────────────────
    right_items = [
        ("Uptake trajectory", f"Y1: {inputs.uptake_y1}%, Y2: {inputs.uptake_y2}%, Y3: {inputs.uptake_y3}%"),
        ("Pricing model", inputs.pricing_model.value),
        ("Price", f"£{inputs.price:,.2f} {inputs.price_unit.value}"),
        ("Setup cost", f"£{inputs.setup_cost:,.2f}"),
        ("Training required", "Yes" if inputs.needs_training else "No"),
    ]

    col_left = Inches(0.6)
    col_right = Inches(6.9)
    col_width = Inches(5.8)
    start_top = Inches(1.35)
    row_height = Inches(1.12)

    for items, left in ((left_items, col_left), (right_items, col_right)):
        for idx, (label, value) in enumerate(items):
            top = start_top + row_height * idx

            # Label
            tb = _add_textbox(slide, left, top, col_width, Inches(0.38))
            run = tb.text_frame.paragraphs[0].add_run()
            run.text = label
            run.font.size = Pt(13)
            run.font.color.rgb = LABEL_GREY
            run.font.name = FONT_NAME

            # Value  (VALUE_WHITE is invisible on white bg — use NAVY)
            tb = _add_textbox(slide, left, top + Inches(0.36), col_width, Inches(0.6))
            run = tb.text_frame.paragraphs[0].add_run()
            run.text = value
            run.font.size = Pt(19)
            run.font.bold = True
            run.font.color.rgb = NAVY
            run.font.name = FONT_NAME


def _add_exec_summary_slide(prs: Presentation, inputs: BIAInputs, results: dict) -> None:
    """Slide 3: Executive Summary with key metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_title_bar(slide, "Executive Summary")

    base = results["base"]
    validation = results.get("validation", {})
    confidence = validation.get("confidence", "N/A")

    # 3-year net impact
    annual_impacts = base["annual_budget_impact"]
    three_year_total = sum(annual_impacts)

    # Break-even
    break_even = base.get("break_even_year")
    break_even_text = f"Year {break_even}" if break_even else "Not reached in forecast period"

    # Top cost drivers
    drivers = base.get("top_cost_drivers", [])
    drivers_text = ", ".join(drivers[:3]) if drivers else "N/A"

    # Layout: 2x2 metric boxes
    metrics = [
        ("3-Year Net Impact", _fmt_gbp(three_year_total), LIGHT_GREEN if three_year_total < 0 else NAVY),
        ("Break-Even Year", break_even_text, GREEN),
        ("Confidence Rating", confidence, NAVY),
        ("Top Cost Drivers", drivers_text, NAVY),
    ]

    positions = [
        (Inches(0.5), Inches(1.35)),
        (Inches(6.85), Inches(1.35)),
        (Inches(0.5), Inches(4.0)),
        (Inches(6.85), Inches(4.0)),
    ]

    for (label, value, accent), (left, top) in zip(metrics, positions):
        # Box background
        _add_shape_rect(slide, left, top, Inches(5.9), Inches(2.45), RGBColor(0xF5, 0xF5, 0xF5))
        # Accent bar at top of box
        _add_shape_rect(slide, left, top, Inches(5.9), Inches(0.09), accent)

        # Label
        tb = _add_textbox(slide, left + Inches(0.3), top + Inches(0.22), Inches(5.3), Inches(0.5))
        tf = tb.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = label
        run.font.size = Pt(16)
        run.font.color.rgb = GREY
        run.font.name = FONT_NAME

        # Value
        tb = _add_textbox(slide, left + Inches(0.3), top + Inches(0.8), Inches(5.3), Inches(1.4))
        tf = tb.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = value
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = accent
        run.font.name = FONT_NAME


def _add_population_slide(prs: Presentation, inputs: BIAInputs) -> None:
    """Slide 5: Population & Uptake details."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_title_bar(slide, "Population & Uptake")

    eligible = inputs.eligible_patients
    treated = inputs.treated_patients_by_year

    # Key stats
    tb = _add_textbox(slide, Inches(0.6), Inches(1.35), Inches(8), Inches(0.65))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = f"Eligible Patients: {eligible:,}"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = GREEN
    run.font.name = FONT_NAME

    tb = _add_textbox(slide, Inches(0.6), Inches(2.1), Inches(11), Inches(0.5))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = (
        f"Catchment: {inputs.catchment_size:,} ({inputs.catchment_type.value})  |  "
        f"Eligible: {inputs.eligible_pct}%"
    )
    run.font.size = Pt(17)
    run.font.color.rgb = GREY
    run.font.name = FONT_NAME

    # Uptake trajectory table
    table = _add_table(
        slide, rows=4, cols=4,
        left=Inches(0.6), top=Inches(2.8),
        width=Inches(12.1), height=Inches(3.8),
    )

    headers = ["", "Year 1", "Year 2", "Year 3"]
    row_data = [
        ["Uptake (%)", f"{inputs.uptake_y1}%", f"{inputs.uptake_y2}%", f"{inputs.uptake_y3}%"],
        ["Treated Patients", f"{treated[0]:,}", f"{treated[1]:,}", f"{treated[2]:,}"],
        ["Cumulative Exposure", f"{treated[0]:,}", f"{treated[0] + treated[1]:,}", f"{sum(treated):,}"],
    ]

    for col_idx, header in enumerate(headers):
        _style_table_cell(table.cell(0, col_idx), header, size_pt=15, bold=True, color=WHITE, bg=NAVY)

    for row_idx, row in enumerate(row_data, start=1):
        for col_idx, val in enumerate(row):
            bold = col_idx == 0
            bg = RGBColor(0xF5, 0xF5, 0xF5) if row_idx % 2 == 0 else None
            _style_table_cell(table.cell(row_idx, col_idx), val, size_pt=14, bold=bold, bg=bg)


def add_budget_impact_table(prs: Presentation, results: BIAResults) -> None:
    """Annual Budget Impact (Base Case) slide with year-by-year table.

    Adds a 4-column table (Year | Treated Patients | Cost per Patient |
    Annual Impact) plus a summary text box showing 3-year net impact and
    break-even year.

    Args:
        prs: The presentation to append the slide to.
        results: A ``BIAResults`` instance (base-case scenario).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_title_bar(slide, "Annual Budget Impact (Base Case)")

    num_years = len(results.annual_budget_impact)
    light_gray = RGBColor(0xF5, 0xF5, 0xF5)

    # ── Table ─────────────────────────────────────────────────────────
    row_h = 0.85
    table = _add_table(
        slide, rows=num_years + 1, cols=4,
        left=Inches(0.6), top=Inches(1.35),
        width=Inches(12.1), height=Inches(row_h + num_years * row_h),
    )

    # Header row
    for col_idx, header in enumerate(
        ["Year", "Treated Patients", "Cost per Patient", "Annual Impact"],
    ):
        _style_table_cell(
            table.cell(0, col_idx), header,
            size_pt=15, bold=True, color=WHITE, bg=NAVY,
            align=PP_ALIGN.CENTER,
        )

    # Data rows
    for yr in range(num_years):
        impact = results.annual_budget_impact[yr]
        row_bg = light_gray if yr % 2 == 1 else None

        _style_table_cell(
            table.cell(yr + 1, 0), f"Year {yr + 1}",
            size_pt=14, bg=row_bg,
        )
        _style_table_cell(
            table.cell(yr + 1, 1),
            f"{results.total_treated_patients[yr]:,}",
            size_pt=14, bg=row_bg, align=PP_ALIGN.RIGHT,
        )
        _style_table_cell(
            table.cell(yr + 1, 2),
            f"£{results.cost_per_patient[yr]:,.2f}",
            size_pt=14, bg=row_bg, align=PP_ALIGN.RIGHT,
        )
        impact_color = LIGHT_GREEN if impact < 0 else RED
        _style_table_cell(
            table.cell(yr + 1, 3),
            f"£{impact:,.0f}",
            size_pt=14, color=impact_color, bg=row_bg,
            align=PP_ALIGN.RIGHT,
        )

    # ── Summary box below the table ───────────────────────────────────
    summary_top = Inches(1.35) + Inches(row_h + num_years * row_h) + Inches(0.45)

    three_year_total = sum(results.annual_budget_impact)
    total_color = LIGHT_GREEN if three_year_total < 0 else RED

    tb = _add_textbox(slide, Inches(0.6), summary_top, Inches(12.1), Inches(0.65))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = f"3-Year Net Impact: £{three_year_total:,.0f}"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = total_color
    run.font.name = FONT_NAME

    be_year = results.break_even_year
    be_text = f"Year {be_year}" if be_year else "Never"
    _add_text(
        tf, f"Break-even: {be_text}",
        size_pt=15, color=GREY, alignment=PP_ALIGN.LEFT,
    )


def add_scenario_comparison(prs: Presentation, scenarios: dict) -> None:
    """Scenario Analysis slide with horizontal bar chart.

    Draws three proportionally-scaled bars (one per scenario) showing
    the 3-year total budget impact, plus an explanatory note.

    Args:
        prs: The presentation to append the slide to.
        scenarios: Dict with keys ``"conservative"``, ``"base"``,
            ``"optimistic"``, each containing at least
            ``"annual_budget_impact"`` (list of floats).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_title_bar(slide, "Scenario Analysis (3-Year Total Impact)")

    order = ["conservative", "base", "optimistic"]
    labels = {"conservative": "Conservative", "base": "Base Case", "optimistic": "Optimistic"}
    bar_colors = {
        "conservative": BAR_CONSERVATIVE,
        "base": BAR_BASE,
        "optimistic": BAR_OPTIMISTIC,
    }

    # ── Calculate 3-year totals ───────────────────────────────────────
    totals = {key: sum(scenarios[key]["annual_budget_impact"]) for key in order}
    max_abs = max(abs(v) for v in totals.values()) or 1

    # ── Layout constants ──────────────────────────────────────────────
    label_left = Inches(0.5)
    label_width = Inches(2.5)
    bar_left = Inches(3.2)
    max_bar_width = 8.5          # inches — longest bar fills this
    bar_height = Inches(0.95)
    start_top = Inches(1.5)
    row_spacing = Inches(1.7)

    for idx, key in enumerate(order):
        total = totals[key]
        top = start_top + row_spacing * idx

        # Scenario label (left)
        tb = _add_textbox(slide, label_left, top + Inches(0.12), label_width, bar_height)
        tf = tb.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        run = tf.paragraphs[0].add_run()
        run.text = labels[key]
        run.font.size = Pt(18)
        run.font.bold = key == "base"
        run.font.color.rgb = NAVY
        run.font.name = FONT_NAME

        # Proportional bar (middle)
        bar_w = max(0.2, abs(total) / max_abs * max_bar_width)
        _add_shape_rect(slide, bar_left, top, Inches(bar_w), bar_height, bar_colors[key])

        # Value label (right of bar)
        value_left = bar_left + Inches(bar_w) + Inches(0.25)
        tb = _add_textbox(slide, value_left, top + Inches(0.06), Inches(3.5), bar_height)
        tf = tb.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = f"£{total:,.0f}"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = bar_colors[key]
        run.font.name = FONT_NAME

    # ── Assumption notes ──────────────────────────────────────────────
    note_top = start_top + row_spacing * 3 + Inches(0.1)
    tb = _add_textbox(slide, Inches(0.5), note_top, Inches(12.3), Inches(0.7))
    tf = tb.text_frame
    tf.paragraphs[0].text = ""

    run = tf.paragraphs[0].add_run()
    run.text = (
        "Conservative: 20% lower uptake, 15% higher costs, 30% lower savings"
    )
    run.font.size = Pt(12)
    run.font.color.rgb = GREY
    run.font.name = FONT_NAME

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = (
        "Optimistic: 20% higher uptake, 10% lower costs, 20% higher savings"
    )
    run.font.size = Pt(12)
    run.font.color.rgb = GREY
    run.font.name = FONT_NAME


def add_assumptions_slide(
    prs: Presentation,
    inputs: BIAInputs,
    top_cost_drivers: list[str],
) -> None:
    """Model Assumptions & Methodology slide.

    Two-column bulleted layout covering perspective, cost sources,
    resource impact assumptions, and top cost drivers.

    Args:
        prs: The presentation to append the slide to.
        inputs: The validated BIAInputs.
        top_cost_drivers: Ranked list of cost driver names (top 3 shown).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_title_bar(slide, "Model Assumptions & Methodology")

    discount_text = (
        "3.5% per NICE guidance"
        if inputs.discounting.value == "on"
        else "None (BIA standard)"
    )

    left_sections = [
        ("Perspective & Timeframe", [
            "NHS payer perspective",
            f"{inputs.forecast_years}-year time horizon",
            f"Model start: FY{inputs.model_year}",
            f"Discounting: {discount_text}",
        ]),
        ("Cost Sources", [
            "Workforce: NHS Agenda for Change 2024/25 rates",
            "Reference costs: NHS National Schedule of Reference Costs",
            "Device pricing: As specified by manufacturer",
        ]),
    ]

    right_sections = [
        ("Resource Impact Assumptions", [
            f"Staff time saved: {inputs.staff_time_saved} mins per patient",
            f"Visit reduction: {inputs.visits_reduced}%",
            f"Complication reduction: {inputs.complications_reduced}%",
            f"Readmission reduction: {inputs.readmissions_reduced}%",
            f"Length of stay reduction: {inputs.los_reduced} days",
        ]),
        ("Top Cost Drivers", list(top_cost_drivers[:3])),
    ]

    for sections, col_left in (
        (left_sections, Inches(0.5)),
        (right_sections, Inches(6.85)),
    ):
        tb = _add_textbox(slide, col_left, Inches(1.3), Inches(5.9), Inches(5.9))
        tf = tb.text_frame
        first_section = True

        for heading, bullets in sections:
            # Section heading
            if first_section:
                p = tf.paragraphs[0]
                first_section = False
            else:
                p = tf.add_paragraph()
                p.space_before = Pt(20)

            run = p.add_run()
            run.text = heading
            run.font.size = Pt(15)
            run.font.bold = True
            run.font.color.rgb = NAVY
            run.font.name = FONT_NAME

            # Bullet items
            for item in bullets:
                p = tf.add_paragraph()
                p.space_before = Pt(4)
                run = p.add_run()
                run.text = f"\u2022  {item}"
                run.font.size = Pt(13)
                run.font.color.rgb = BLACK
                run.font.name = FONT_NAME


# ── Main entry point ──────────────────────────────────────────────────

def generate_bia_report(
    inputs: BIAInputs,
    results: dict,
    submission_id: str,
) -> str:
    """Generate a branded PPTX slide deck from BIA inputs and results.

    Args:
        inputs: The validated BIAInputs used for this analysis.
        results: The full API response dict containing ``base``,
            ``conservative``, ``optimistic`` (each a dict with
            ``annual_budget_impact``, ``cost_per_patient``, etc.),
            plus ``validation`` and ``summary``.
        submission_id: Used for the output filename.

    Returns:
        Absolute filepath of the saved ``.pptx`` file.
    """
    prs = Presentation()
    set_heor_theme(prs)

    report_date = date.today().strftime("%d %B %Y")
    intervention_name = f"Budget Impact Analysis — {inputs.setting.value}"

    # Slide 1: Title
    add_title_slide(prs, intervention_name, report_date)

    # Slide 2: Input Summary
    add_input_summary_slide(prs, inputs)

    # Slide 3: Section divider — Executive Summary
    add_section_divider(prs, "Executive Summary")

    # Slide 4: Executive Summary content
    _add_exec_summary_slide(prs, inputs, results)

    # Slide 5: Section divider — Population & Uptake
    add_section_divider(prs, "Population & Uptake")

    # Slide 6: Population & Uptake content
    _add_population_slide(prs, inputs)

    # Slide 7: Section divider — Budget Impact
    add_section_divider(prs, "Budget Impact")

    # Slide 8: Annual Budget Impact (base case)
    base_results = BIAResults(**results["base"])
    add_budget_impact_table(prs, base_results)

    # Slide 9: Scenario Comparison
    add_scenario_comparison(prs, results)

    # Slide 10: Model Assumptions & Methodology
    add_assumptions_slide(prs, inputs, results["base"].get("top_cost_drivers", []))

    # Save
    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    filepath = REPORTS_DIR / f"{submission_id}.pptx"
    prs.save(str(filepath))

    return str(filepath)


# =====================================================================
# CEA (Markov) slide builders
# =====================================================================

def _add_cea_title_slide(
    prs: Presentation,
    intervention_name: str,
    report_date: str,
) -> None:
    """CEA title slide — navy background, HEOR Engine branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_shape_rect(slide, Inches(0), Inches(0), Emu(SLIDE_WIDTH), Emu(SLIDE_HEIGHT), NAVY)
    _add_shape_rect(slide, Inches(0.6), Inches(1.5), Inches(1.5), Inches(0.06), LIGHT_GREEN)

    # Branding
    tb = _add_textbox(slide, Inches(0.6), Inches(0.6), Inches(8), Inches(0.8))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "HEOR Engine"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = LIGHT_GREEN
    run.font.name = FONT_NAME

    # Intervention name
    tb = _add_textbox(slide, Inches(0.6), Inches(1.8), Inches(11), Inches(1.5))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = intervention_name
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    # Subtitle
    tb = _add_textbox(slide, Inches(0.6), Inches(3.5), Inches(11), Inches(1))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "Cost-Effectiveness Analysis"
    run.font.size = Pt(22)
    run.font.color.rgb = LIGHT_GREEN
    run.font.name = FONT_NAME
    _add_text(tf, report_date, size_pt=16, color=WHITE, alignment=PP_ALIGN.LEFT)


def _add_model_structure_slide(
    prs: Presentation,
    inputs: MarkovInputs,
) -> None:
    """Markov Model Structure — state diagram and parameter summary."""
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_title_bar(slide, "Markov Model Structure")

    # ── State boxes ──────────────────────────────────────────────────
    alive_left, alive_top = Inches(1.0), Inches(1.8)
    dead_left, dead_top = Inches(8.5), Inches(1.8)
    box_w, box_h = Inches(3.2), Inches(1.5)

    # Alive box (green)
    alive_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, alive_left, alive_top, box_w, box_h,
    )
    alive_box.fill.solid()
    alive_box.fill.fore_color.rgb = GREEN
    alive_box.line.fill.background()
    alive_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = alive_box.text_frame.paragraphs[0].add_run()
    run.text = "Alive"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    # Dead box (dark grey)
    dead_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, dead_left, dead_top, box_w, box_h,
    )
    dead_box.fill.solid()
    dead_box.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x44)
    dead_box.line.fill.background()
    dead_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = dead_box.text_frame.paragraphs[0].add_run()
    run.text = "Dead"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    # Arrow (rectangle between boxes)
    arrow_left = alive_left + box_w + Inches(0.2)
    arrow_top = alive_top + box_h / 2 - Inches(0.03)
    arrow_w = dead_left - arrow_left - Inches(0.2)
    _add_shape_rect(slide, arrow_left, arrow_top, arrow_w, Inches(0.06), LIGHT_GREEN)

    # Arrow head
    head = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, dead_left - Inches(0.3),
        alive_top + box_h / 2 - Inches(0.15), Inches(0.3), Inches(0.3),
    )
    head.fill.solid()
    head.fill.fore_color.rgb = LIGHT_GREEN
    head.line.fill.background()
    head.rotation = 90.0

    # Transition label
    tb = _add_textbox(
        slide,
        arrow_left + Inches(0.2), arrow_top - Inches(0.55),
        Inches(3), Inches(0.45),
    )
    tf = tb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = f"Standard: {inputs.prob_death_standard:.0%}  |  Treatment: {inputs.prob_death_treatment:.0%}"
    run.font.size = Pt(13)
    run.font.color.rgb = NAVY
    run.font.name = FONT_NAME

    # Self-loop label on Alive
    tb = _add_textbox(
        slide,
        alive_left, alive_top + box_h + Inches(0.12),
        box_w, Inches(0.4),
    )
    tf = tb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = "Costs & QALYs accrue"
    run.font.size = Pt(13)
    run.font.color.rgb = GREY
    run.font.name = FONT_NAME

    # ── Parameter summary table ──────────────────────────────────────
    table = _add_table(
        slide, rows=4, cols=2,
        left=Inches(1.5), top=Inches(4.2),
        width=Inches(9.8), height=Inches(2.7),
    )
    params = [
        ("Time horizon", f"{inputs.time_horizon} years"),
        ("Cycle length", f"{inputs.cycle_length} year" if inputs.cycle_length == 1 else f"{inputs.cycle_length} years (quarterly)" if inputs.cycle_length == 0.25 else f"{inputs.cycle_length}"),
        ("Discount rate", f"{inputs.discount_rate:.1%} (NICE standard)" if inputs.discount_rate == 0.035 else f"{inputs.discount_rate:.1%}"),
    ]
    _style_table_cell(table.cell(0, 0), "Parameter", size_pt=14, bold=True, color=WHITE, bg=NAVY)
    _style_table_cell(table.cell(0, 1), "Value", size_pt=14, bold=True, color=WHITE, bg=NAVY)
    for i, (param, val) in enumerate(params, 1):
        bg = LIGHT_GREY if i % 2 == 0 else None
        _style_table_cell(table.cell(i, 0), param, size_pt=13, bg=bg)
        _style_table_cell(table.cell(i, 1), val, size_pt=13, bold=True, bg=bg)


def _add_cea_inputs_slide(prs: Presentation, inputs: MarkovInputs) -> None:
    """Input Parameters slide — two-column comparison table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_title_bar(slide, "Input Parameters")

    table = _add_table(
        slide, rows=5, cols=3,
        left=Inches(0.5), top=Inches(1.35),
        width=Inches(12.3), height=Inches(4.8),
    )

    # Header
    _style_table_cell(table.cell(0, 0), "Parameter", size_pt=15, bold=True, color=WHITE, bg=NAVY)
    _style_table_cell(table.cell(0, 1), "Standard Care", size_pt=15, bold=True, color=WHITE, bg=NAVY, align=PP_ALIGN.CENTER)
    _style_table_cell(table.cell(0, 2), "Treatment", size_pt=15, bold=True, color=WHITE, bg=GREEN, align=PP_ALIGN.CENTER)

    rows = [
        ("Annual mortality", f"{inputs.prob_death_standard:.1%}", f"{inputs.prob_death_treatment:.1%}"),
        ("Annual cost", f"£{inputs.cost_standard_annual:,.0f}", f"£{inputs.cost_treatment_annual:,.0f}"),
        ("Initial cost", "—", f"£{inputs.cost_treatment_initial:,.0f}" if inputs.cost_treatment_initial else "—"),
        ("Quality of life (utility)", f"{inputs.utility_standard:.2f}", f"{inputs.utility_treatment:.2f}"),
    ]
    for i, (param, std, trt) in enumerate(rows, 1):
        bg = LIGHT_GREY if i % 2 == 0 else None
        _style_table_cell(table.cell(i, 0), param, size_pt=14, bold=True, bg=bg)
        _style_table_cell(table.cell(i, 1), std, size_pt=14, bg=bg, align=PP_ALIGN.CENTER)
        _style_table_cell(table.cell(i, 2), trt, size_pt=14, bg=bg, align=PP_ALIGN.CENTER)

    # Footnote
    tb = _add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = (
        f"Discount rate: {inputs.discount_rate:.1%}  |  "
        f"Time horizon: {inputs.time_horizon} years  |  "
        f"Cycle length: {'Annual' if inputs.cycle_length == 1 else f'{inputs.cycle_length} year'}"
    )
    run.font.size = Pt(12)
    run.font.color.rgb = GREY
    run.font.name = FONT_NAME


def _add_cea_results_slide(prs: Presentation, results: MarkovResults) -> None:
    """Results slide — costs, QALYs, and ICER callout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_title_bar(slide, "Cost-Effectiveness Results")

    # ── Results table ────────────────────────────────────────────────
    table = _add_table(
        slide, rows=4, cols=4,
        left=Inches(0.5), top=Inches(1.35),
        width=Inches(12.3), height=Inches(2.8),
    )

    headers = ["", "Total Cost", "Total QALYs", "Cost per QALY"]
    for i, h in enumerate(headers):
        _style_table_cell(table.cell(0, i), h, size_pt=15, bold=True, color=WHITE, bg=NAVY, align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    std = results.standard_care
    trt = results.treatment

    # Standard care row
    _style_table_cell(table.cell(1, 0), "Standard Care", size_pt=14, bold=True)
    _style_table_cell(table.cell(1, 1), _fmt_gbp(std.total_cost), size_pt=14, align=PP_ALIGN.RIGHT)
    _style_table_cell(table.cell(1, 2), f"{std.total_qalys:.4f}", size_pt=14, align=PP_ALIGN.RIGHT)
    _style_table_cell(table.cell(1, 3), "—", size_pt=14, color=GREY, align=PP_ALIGN.CENTER)

    # Treatment row
    _style_table_cell(table.cell(2, 0), "Treatment", size_pt=14, bold=True, bg=LIGHT_GREY)
    _style_table_cell(table.cell(2, 1), _fmt_gbp(trt.total_cost), size_pt=14, bg=LIGHT_GREY, align=PP_ALIGN.RIGHT)
    _style_table_cell(table.cell(2, 2), f"{trt.total_qalys:.4f}", size_pt=14, bg=LIGHT_GREY, align=PP_ALIGN.RIGHT)
    _style_table_cell(table.cell(2, 3), "—", size_pt=14, color=GREY, bg=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # Incremental row
    cost_color = LIGHT_GREEN if results.incremental_cost < 0 else RED
    qaly_color = LIGHT_GREEN if results.incremental_qalys > 0 else RED
    icer_text = _fmt_gbp(results.icer) if results.icer is not None else "N/A"

    _style_table_cell(table.cell(3, 0), "Incremental", size_pt=15, bold=True, color=WHITE, bg=GREEN)
    _style_table_cell(table.cell(3, 1), _fmt_gbp(results.incremental_cost), size_pt=15, bold=True, color=cost_color, bg=RGBColor(0xE8, 0xF5, 0xE9), align=PP_ALIGN.RIGHT)
    _style_table_cell(table.cell(3, 2), f"{results.incremental_qalys:.4f}", size_pt=15, bold=True, color=qaly_color, bg=RGBColor(0xE8, 0xF5, 0xE9), align=PP_ALIGN.RIGHT)
    _style_table_cell(table.cell(3, 3), icer_text, size_pt=15, bold=True, color=WHITE, bg=GREEN, align=PP_ALIGN.CENTER)

    # ── ICER callout box ─────────────────────────────────────────────
    callout_top = Inches(4.45)
    badge_color = LIGHT_GREEN if results.cost_effective_25k else AMBER if results.cost_effective_35k else RED

    _add_shape_rect(slide, Inches(0.9), callout_top, Inches(11.5), Inches(2.4), RGBColor(0xF5, 0xF8, 0xFC))
    _add_shape_rect(slide, Inches(0.9), callout_top, Inches(0.1), Inches(2.4), badge_color)

    # ICER label
    tb = _add_textbox(slide, Inches(1.3), callout_top + Inches(0.2), Inches(10.8), Inches(0.38))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "INCREMENTAL COST-EFFECTIVENESS RATIO"
    run.font.size = Pt(13)
    run.font.color.rgb = GREY
    run.font.name = FONT_NAME

    # ICER value
    tb = _add_textbox(slide, Inches(1.3), callout_top + Inches(0.62), Inches(10.8), Inches(1.1))
    run = tb.text_frame.paragraphs[0].add_run()
    icer_display = f"{_fmt_gbp(results.icer)} per QALY" if results.icer is not None else "N/A — no QALY difference"
    run.text = icer_display
    run.font.size = Pt(38)
    run.font.bold = True
    run.font.color.rgb = badge_color
    run.font.name = FONT_NAME

    # Interpretation
    tb = _add_textbox(slide, Inches(1.3), callout_top + Inches(1.75), Inches(10.8), Inches(0.5))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = results.interpretation
    run.font.size = Pt(16)
    run.font.color.rgb = NAVY
    run.font.name = FONT_NAME


def _add_ce_plane_slide(prs: Presentation, results: MarkovResults) -> None:
    """Cost-Effectiveness Plane — scatter plot built from shapes.

    Plots the intervention as a point on the incremental cost vs.
    incremental QALYs plane, with quadrant labels and NICE threshold lines.
    """
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_title_bar(slide, "Cost-Effectiveness Plane")

    # ── Plot area geometry ───────────────────────────────────────────
    plot_left = Inches(0.8)
    plot_top = Inches(1.35)
    plot_w = Inches(11.7)
    plot_h = Inches(5.7)

    # White background for plot
    _add_shape_rect(slide, plot_left, plot_top, plot_w, plot_h, WHITE)

    # Border
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, plot_left, plot_top, plot_w, plot_h,
    )
    border.fill.background()
    border.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    border.line.width = Pt(1)

    # ── Axes (cross at centre) ───────────────────────────────────────
    centre_x = plot_left + plot_w // 2
    centre_y = plot_top + plot_h // 2

    # Horizontal axis (QALYs)
    _add_shape_rect(slide, plot_left, centre_y - Inches(0.01), plot_w, Inches(0.02), NAVY)
    # Vertical axis (Cost)
    _add_shape_rect(slide, centre_x - Inches(0.01), plot_top, Inches(0.02), plot_h, NAVY)

    # ── Axis labels ──────────────────────────────────────────────────
    tb = _add_textbox(slide, centre_x + Inches(2.5), centre_y + Inches(0.1), Inches(2.5), Inches(0.4))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Incremental QALYs \u2192"
    run.font.size = Pt(12)
    run.font.color.rgb = NAVY
    run.font.name = FONT_NAME

    tb = _add_textbox(slide, plot_left, plot_top - Inches(0.3), Inches(3.5), Inches(0.35))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "\u2191 Incremental Cost (\u00A3)"
    run.font.size = Pt(12)
    run.font.color.rgb = NAVY
    run.font.name = FONT_NAME

    # ── Quadrant labels ──────────────────────────────────────────────
    quadrants = [
        ("NW: Dominated\n(More costly, less effective)", plot_left + Inches(0.3), plot_top + Inches(0.2)),
        ("NE: Trade-off\n(More costly, more effective)", centre_x + Inches(0.5), plot_top + Inches(0.2)),
        ("SW: Trade-off\n(Less costly, less effective)", plot_left + Inches(0.3), centre_y + Inches(0.2)),
        ("SE: Dominant\n(Less costly, more effective)", centre_x + Inches(0.5), centre_y + Inches(0.2)),
    ]
    for text, lft, tp in quadrants:
        tb = _add_textbox(slide, lft, tp, Inches(3.5), Inches(0.7))
        for line_idx, line in enumerate(text.split("\n")):
            if line_idx == 0:
                run = tb.text_frame.paragraphs[0].add_run()
            else:
                p = tb.text_frame.add_paragraph()
                run = p.add_run()
            run.text = line
            run.font.size = Pt(11)
            run.font.color.rgb = GREY
            run.font.name = FONT_NAME

    # ── NICE threshold lines (£20k and £30k per QALY) ───────────────
    for threshold, line_color, dash_label in [
        (20_000, LIGHT_GREEN, "£20k/QALY"),
        (30_000, AMBER, "£30k/QALY"),
    ]:
        label_x = centre_x + Inches(2.8)
        label_y_offset = Inches(1.4) if threshold == 20_000 else Inches(0.7)
        label_y = centre_y - label_y_offset

        _add_shape_rect(slide, label_x - Inches(0.7), label_y + Inches(0.1), Inches(0.6), Inches(0.05), line_color)

        tb = _add_textbox(slide, label_x, label_y, Inches(2), Inches(0.38))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = dash_label
        run.font.size = Pt(11)
        run.font.color.rgb = line_color
        run.font.name = FONT_NAME

    # ── Plot the intervention point ──────────────────────────────────
    inc_q = results.incremental_qalys
    inc_c = results.incremental_cost

    q_range = max(abs(inc_q) * 2.5, 0.5)
    c_range = max(abs(inc_c) * 2.5, 10000)

    px = centre_x + int(inc_q / q_range * (plot_w // 2))
    py = centre_y - int(inc_c / c_range * (plot_h // 2))

    dot_size = Inches(0.32)
    px = max(plot_left, min(px, plot_left + plot_w - dot_size))
    py = max(plot_top, min(py, plot_top + plot_h - dot_size))

    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, px - dot_size // 2, py - dot_size // 2,
        dot_size, dot_size,
    )
    dot.fill.solid()
    dot_color = LIGHT_GREEN if results.cost_effective_25k else AMBER if results.cost_effective_35k else RED
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()

    # Label next to dot
    tb = _add_textbox(slide, px + Inches(0.25), py - Inches(0.35), Inches(3.5), Inches(0.6))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    icer_text = _fmt_gbp(results.icer) + "/QALY" if results.icer is not None else "N/A"
    run.text = f"ICER: {icer_text}"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = dot_color
    run.font.name = FONT_NAME

    _add_text(tf, f"({_fmt_gbp(inc_c)}, {inc_q:.4f} QALYs)", size_pt=11, color=GREY)


def _add_interpretation_slide(prs: Presentation, inputs: MarkovInputs, results: MarkovResults) -> None:
    """Interpretation slide — large verdict, threshold indicators, assumptions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_navy_title_bar(slide, "Interpretation & Recommendations")

    # ── Main verdict box ─────────────────────────────────────────────
    verdict_color = LIGHT_GREEN if results.cost_effective_25k else AMBER if results.cost_effective_35k else RED

    _add_shape_rect(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(2.0), RGBColor(0xF5, 0xF8, 0xFC))
    _add_shape_rect(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.1), verdict_color)

    tb = _add_textbox(slide, Inches(0.9), Inches(1.55), Inches(11.5), Inches(1.5))
    tf = tb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = results.interpretation
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = verdict_color
    run.font.name = FONT_NAME

    icer_summary = f"ICER: {_fmt_gbp(results.icer)}/QALY" if results.icer is not None else "No QALY difference between arms"
    _add_text(tf, icer_summary, size_pt=17, color=GREY, alignment=PP_ALIGN.CENTER)

    # ── NICE threshold checklist ─────────────────────────────────────
    thresholds = [
        ("\u00A325,000/QALY — Standard threshold", results.cost_effective_25k),
        ("\u00A335,000/QALY — Extended threshold", results.cost_effective_35k),
        ("\u00A350,000/QALY — End-of-life criteria", results.icer is not None and results.icer < 50_000),
    ]

    tb = _add_textbox(slide, Inches(0.5), Inches(3.6), Inches(6.0), Inches(3.2))
    tf = tb.text_frame

    run = tf.paragraphs[0].add_run()
    run.text = "NICE Threshold Assessment"
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run.font.name = FONT_NAME

    for label, passed in thresholds:
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        icon = "\u2713" if passed else "\u2717"
        color = GREEN if passed else RED
        run = p.add_run()
        run.text = f"  {icon}  {label}"
        run.font.size = Pt(15)
        run.font.color.rgb = color
        run.font.name = FONT_NAME

    # ── Key assumptions ──────────────────────────────────────────────
    tb = _add_textbox(slide, Inches(7.0), Inches(3.6), Inches(5.8), Inches(3.2))
    tf = tb.text_frame

    run = tf.paragraphs[0].add_run()
    run.text = "Key Assumptions"
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run.font.name = FONT_NAME

    assumptions = [
        "2-state Markov model (Alive / Dead)",
        f"{inputs.time_horizon}-year time horizon, {'annual' if inputs.cycle_length == 1 else f'{inputs.cycle_length}-year'} cycles",
        f"Discount rate: {inputs.discount_rate:.1%}",
        "Constant transition probabilities over time",
        "Beginning-of-cycle cost & QALY accrual",
        "NHS payer perspective",
    ]
    for a in assumptions:
        p = tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"\u2022  {a}"
        run.font.size = Pt(13)
        run.font.color.rgb = BLACK
        run.font.name = FONT_NAME


# ── CEA public entry points ──────────────────────────────────────────

def generate_cea_report(
    inputs: MarkovInputs,
    results: MarkovResults,
    submission_id: str,
) -> str:
    """Generate a branded PPTX slide deck for a Markov CEA.

    Creates a 6-slide presentation:
        1. Title slide
        2. Markov model structure diagram
        3. Input parameters comparison table
        4. Results — costs, QALYs, ICER callout
        5. Cost-effectiveness plane
        6. Interpretation & recommendations

    Args:
        inputs: Validated ``MarkovInputs``.
        results: ``MarkovResults`` from the R model run.
        submission_id: Used for the output filename.

    Returns:
        Absolute filepath of the saved ``.pptx`` file.
    """
    prs = Presentation()
    set_heor_theme(prs)

    report_date = date.today().strftime("%d %B %Y")

    _add_cea_title_slide(prs, inputs.intervention_name, report_date)
    _add_model_structure_slide(prs, inputs)
    _add_cea_inputs_slide(prs, inputs)
    _add_cea_results_slide(prs, results)
    _add_ce_plane_slide(prs, results)
    _add_interpretation_slide(prs, inputs, results)

    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    filepath = REPORTS_DIR / f"CEA_{submission_id}.pptx"
    prs.save(str(filepath))

    return str(filepath)


def add_cea_slides_to_bia_report(
    bia_inputs: BIAInputs,
    bia_results: dict,
    markov_inputs: MarkovInputs,
    markov_results: MarkovResults,
    submission_id: str,
) -> str:
    """Generate a combined BIA + CEA report.

    Builds the standard 10-slide BIA deck and appends 5 CEA slides
    (model structure, inputs, results, CE plane, interpretation)
    after a section divider, for a total of 16 slides.

    Args:
        bia_inputs: Validated BIA inputs.
        bia_results: Full BIA API response dict.
        markov_inputs: Validated Markov inputs.
        markov_results: Markov results from the R model.
        submission_id: Used for the output filename.

    Returns:
        Absolute filepath of the saved ``.pptx`` file.
    """
    prs = Presentation()
    set_heor_theme(prs)

    report_date = date.today().strftime("%d %B %Y")
    intervention_name = f"Budget Impact & Cost-Effectiveness — {bia_inputs.setting.value}"

    # ── BIA slides (1–10) ────────────────────────────────────────────
    add_title_slide(prs, intervention_name, report_date)
    add_input_summary_slide(prs, bia_inputs)
    add_section_divider(prs, "Executive Summary")
    _add_exec_summary_slide(prs, bia_inputs, bia_results)
    add_section_divider(prs, "Population & Uptake")
    _add_population_slide(prs, bia_inputs)
    add_section_divider(prs, "Budget Impact")
    base_results = BIAResults(**bia_results["base"])
    add_budget_impact_table(prs, base_results)
    add_scenario_comparison(prs, bia_results)
    add_assumptions_slide(prs, bia_inputs, bia_results["base"].get("top_cost_drivers", []))

    # ── CEA section divider + slides (11–16) ─────────────────────────
    add_section_divider(prs, "Cost-Effectiveness Analysis")
    _add_model_structure_slide(prs, markov_inputs)
    _add_cea_inputs_slide(prs, markov_inputs)
    _add_cea_results_slide(prs, markov_results)
    _add_ce_plane_slide(prs, markov_results)
    _add_interpretation_slide(prs, markov_inputs, markov_results)

    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    filepath = REPORTS_DIR / f"BIA_CEA_{submission_id}.pptx"
    prs.save(str(filepath))

    return str(filepath)
